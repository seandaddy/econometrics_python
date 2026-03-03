#!/usr/bin/env python
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re

def parse_markdown(md_file):
    """Parse markdown file and extract slides"""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Split by slide separator (---)
    slides = content.split('\n---\n')
    
    parsed_slides = []
    for slide in slides:
        slide = slide.strip()
        if not slide:
            continue
        
        lines = slide.split('\n')
        slide_data = {
            'title': '',
            'content': []
        }
        
        i = 0
        while i < len(lines):
            line = lines[i].strip()
            
            # H1 title
            if line.startswith('# '):
                slide_data['title'] = line[2:].strip()
                i += 1
            # H2 title (if no H1)
            elif line.startswith('## ') and not slide_data['title']:
                slide_data['title'] = line[3:].strip()
                i += 1
            # Skip empty lines
            elif not line:
                i += 1
            else:
                break
        
        # Collect remaining content
        content_lines = []
        current_section = None
        
        while i < len(lines):
            line = lines[i].strip()
            
            if line.startswith('## '):
                # Subtitle
                if current_section:
                    content_lines.append(current_section)
                current_section = {'type': 'subtitle', 'text': line[3:].strip()}
            elif line.startswith('### '):
                # Section header
                if current_section:
                    content_lines.append(current_section)
                current_section = {'type': 'header', 'text': line[4:].strip()}
            elif line.startswith('#### '):
                # Sub-header
                content_lines.append({'type': 'subheader', 'text': line[5:].strip()})
            elif line.startswith('- '):
                # Bullet point
                content_lines.append({'type': 'bullet', 'text': line[2:].strip()})
            elif line.startswith('```'):
                # Code block
                code_lines = []
                i += 1
                while i < len(lines) and not lines[i].strip().startswith('```'):
                    code_lines.append(lines[i])
                    i += 1
                content_lines.append({'type': 'code', 'text': '\n'.join(code_lines)})
            elif re.match(r'^\d+\.\s', line):
                # Numbered list
                content_lines.append({'type': 'numbered', 'text': line})
            elif line.startswith('✓ '):
                # Checkmark bullet
                content_lines.append({'type': 'checkmark', 'text': line[2:].strip()})
            elif line:
                # Regular text
                content_lines.append({'type': 'text', 'text': line})
            
            i += 1
        
        if current_section:
            content_lines.append(current_section)
        
        slide_data['content'] = content_lines
        parsed_slides.append(slide_data)
    
    return parsed_slides

def create_presentation(template_file, parsed_slides, output_file):
    """Create PowerPoint presentation using template"""
    
    try:
        # Try to use template
        prs = Presentation(template_file)
        print(f"Using template: {template_file}")
    except:
        # Create new presentation if template fails
        prs = Presentation()
        print("Creating new presentation (template not accessible)")
    
    # Set slide size (16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Check available layouts
    num_layouts = len(prs.slide_layouts)
    print(f"Available layouts: {num_layouts}")
    
    # Use the last available layout (usually blank or most flexible)
    layout_index = min(num_layouts - 1, 6) if num_layouts > 0 else 0
    
    for slide_data in parsed_slides:
        # Add a blank slide
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        if slide_data['title']:
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
            )
            title_frame = title_box.text_frame
            title_frame.text = slide_data['title']
            
            # Format title
            for paragraph in title_frame.paragraphs:
                paragraph.font.size = Pt(32)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(0, 51, 102)
                paragraph.alignment = PP_ALIGN.LEFT
        
        # Add content
        top = 1.3
        left = 0.7
        width = 8.6
        
        content_box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(5.5)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for item in slide_data['content']:
            p = text_frame.add_paragraph()
            
            if item['type'] == 'subtitle':
                p.text = item['text']
                p.font.size = Pt(24)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 102, 204)
                p.space_before = Pt(12)
                p.space_after = Pt(8)
            
            elif item['type'] == 'header':
                p.text = item['text']
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = RGBColor(51, 51, 51)
                p.space_before = Pt(10)
                p.space_after = Pt(6)
            
            elif item['type'] == 'subheader':
                p.text = item['text']
                p.font.size = Pt(18)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 102, 204)
                p.space_before = Pt(8)
                p.space_after = Pt(4)
            
            elif item['type'] == 'bullet':
                p.text = item['text']
                p.font.size = Pt(16)
                p.level = 0
                p.space_after = Pt(4)
            
            elif item['type'] == 'numbered':
                p.text = item['text']
                p.font.size = Pt(16)
                p.space_after = Pt(4)
            
            elif item['type'] == 'checkmark':
                p.text = '✓ ' + item['text']
                p.font.size = Pt(16)
                p.space_after = Pt(4)
            
            elif item['type'] == 'code':
                p.text = item['text']
                p.font.name = 'Courier New'
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(51, 51, 51)
                p.space_before = Pt(6)
                p.space_after = Pt(6)
            
            elif item['type'] == 'text':
                p.text = item['text']
                p.font.size = Pt(16)
                p.space_after = Pt(4)
    
    # Save presentation
    prs.save(output_file)
    print(f"✓ Presentation saved: {output_file}")
    print(f"✓ Total slides: {len(parsed_slides)}")

if __name__ == "__main__":
    md_file = "chapter1_presentation.md"
    template_file = "class101_template.pptx"
    output_file = "chapter1_presentation.pptx"
    
    print("Parsing markdown file...")
    parsed_slides = parse_markdown(md_file)
    
    print(f"Creating presentation with {len(parsed_slides)} slides...")
    create_presentation(template_file, parsed_slides, output_file)
    
    print("\n✓ Done!")
